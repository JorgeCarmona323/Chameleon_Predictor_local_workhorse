# env: chameleon-calc
"""
export_slides.py
----------------
Renders the Chameleon Predictor methodology flowchart as a PNG and
inserts it into a PowerPoint slide alongside the key result figures.

Usage:
    python scripts/export_slides.py [--outdir results]

Outputs:
    results/methodology_flowchart.png
    results/Chameleon_Predictor_slides.pptx
"""

import argparse
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colours ───────────────────────────────────────────────────────────────────
C_DATA    = "#4A90D9"   # blue   — data / input nodes
C_SCRIPT  = "#7B68EE"  # purple — pipeline scripts
C_RESULT  = "#2ECC71"  # green  — good result
C_WARN    = "#F39C12"  # amber  — noisy / heterogeneous result
C_NEG     = "#E74C3C"  # red    — negative control
C_NEUTRAL = "#95A5A6"  # grey   — intermediate / output
C_BG      = "#FAFAFA"  # near-white background
WHITE     = "#FFFFFF"


def draw_box(ax, x, y, w, h, label, color, fontsize=7.5, text_color="white",
             style="round,pad=0.1", zorder=3):
    box = FancyBboxPatch((x - w/2, y - h/2), w, h,
                         boxstyle=style, linewidth=0.8,
                         edgecolor="#333333", facecolor=color, zorder=zorder)
    ax.add_patch(box)
    ax.text(x, y, label, ha="center", va="center", fontsize=fontsize,
            color=text_color, zorder=zorder+1,
            multialignment="center", linespacing=1.35)


def arrow(ax, x0, y0, x1, y1, color="#555555"):
    ax.annotate("", xy=(x1, y1), xytext=(x0, y0),
                arrowprops=dict(arrowstyle="-|>", color=color,
                                lw=1.0, mutation_scale=10))


def make_flowchart(out_path: Path):
    fig, ax = plt.subplots(figsize=(14, 9))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 9)
    ax.axis("off")
    fig.patch.set_facecolor(C_BG)
    ax.set_facecolor(C_BG)

    # ── Row positions (top → bottom) ──────────────────────────────────────────
    R = {
        "db":      8.4,
        "curate":  7.5,
        "subset":  6.6,
        "feat":    5.5,
        "matrix":  4.4,
        "analysis":3.3,
        "umap":    2.2,
        "tracks":  1.1,
        "strat":   0.35,
    }

    BW, BH = 2.6, 0.58   # default box width / height
    SW, SH = 1.9, 0.52   # small box
    TW, TH = 1.3, 0.48   # track box

    # ── Database ──────────────────────────────────────────────────────────────
    draw_box(ax, 7, R["db"], BW, BH,
             "CycPeptMPDB v1.2\n8,466 cyclic peptides", C_DATA)

    arrow(ax, 7, R["db"]-BH/2, 7, R["curate"]+BH/2)

    # ── Curation ──────────────────────────────────────────────────────────────
    draw_box(ax, 7, R["curate"], BW, BH,
             "curate_data.py\nPAMPA filter · RDKit canonicalization", C_SCRIPT)

    arrow(ax, 7, R["curate"]-BH/2, 7, R["subset"]+BH/2)

    # ── PAMPA subset ──────────────────────────────────────────────────────────
    draw_box(ax, 7, R["subset"], BW, BH,
             "PAMPA subset\n7,298 compounds", C_DATA)

    # Three branches from subset
    xs_feat = [2.5, 7, 11.5]
    labels_feat = [
        "2D Descriptors\nMolWt · MolLogP · TPSA\nHBA · HBD · RotBonds",
        "DB 3DPSA\nH₂O_3DPSA − CHCl₃_3DPSA\n(negative control)",
        "Tier-1: ETKDGv3 + MMFF94s\n20 conformers / molecule\nΔPSA · ΔHB · ΔRg · PSA_std",
    ]
    colors_feat = [C_NEUTRAL, C_NEG, C_SCRIPT]

    for xf, lf, cf in zip(xs_feat, labels_feat, colors_feat):
        arrow(ax, 7, R["subset"]-BH/2,
              xf, R["feat"]+SH/2)
        draw_box(ax, xf, R["feat"], SW, SH*1.3, lf, cf, fontsize=6.8)

    # Merge arrows to feature matrix
    for xf in xs_feat:
        arrow(ax, xf, R["feat"]-SH*0.65, 7, R["matrix"]+BH/2)

    # ── Feature matrix ────────────────────────────────────────────────────────
    draw_box(ax, 7, R["matrix"], BW, BH,
             "build_feature_matrix.py\nmerge all features · 7,298 rows", C_SCRIPT)

    # Split to analysis and UMAP
    arrow(ax, 4.5, R["matrix"], 3.5, R["analysis"]+BH/2)
    arrow(ax, 9.5, R["matrix"], 10.5, R["analysis"]+BH/2)
    arrow(ax, 7, R["matrix"]-BH/2, 7, R["umap"]+BH/2)

    # ── Parallel: correlation and UMAP ────────────────────────────────────────
    draw_box(ax, 2.8, R["analysis"], 2.8, BH,
             "correlation_analysis.py\nPearson · Spearman · AUC-ROC", C_SCRIPT,
             fontsize=6.8)

    draw_box(ax, 7, R["umap"], BW, BH,
             "umap_visualization.py\nRobustScaler → UMAP cosine", C_SCRIPT)

    draw_box(ax, 11.2, R["analysis"], 2.8, BH,
             "Enrichment tables\nper-cluster perm rate\ndouble-validated islands",
             C_NEUTRAL, fontsize=6.8)

    # ── Four UMAP tracks ─────────────────────────────────────────────────────
    track_xs = [3.0, 5.5, 8.5, 11.0]
    track_labels = [
        "Track A\nK-Medoids\narchetypes",
        "Track B\nHDBSCAN\nclusters",
        "Track C\nPAMPA\nLogPexp",
        "Track D\nMol Weight\nMW coloring",
    ]
    track_colors = [C_SCRIPT, C_SCRIPT, C_DATA, C_SCRIPT]

    for xt, lt, ct in zip(track_xs, track_labels, track_colors):
        arrow(ax, 7, R["umap"]-BH/2, xt, R["tracks"]+TH/2)
        draw_box(ax, xt, R["tracks"], TW, TH, lt, ct, fontsize=6.5)

    # ARI stability from Track A+B
    arrow(ax, 4.25, R["tracks"]-TH/2, 4.25, R["strat"]+0.12)
    draw_box(ax, 4.25, R["strat"]-0.02, 1.8, 0.36,
             "ARI stability\n5 seeds · threshold ≥0.85",
             C_NEUTRAL, fontsize=6.2)

    # ── Source stratification ─────────────────────────────────────────────────
    arrow(ax, 7, R["matrix"]-BH/2-0.05, 9.5, R["strat"]+0.12)
    draw_box(ax, 9.6, R["strat"]+0.13, 1.7, 0.36,
             "Source filter\n--sources flag", C_SCRIPT, fontsize=6.2)

    draw_box(ax, 8.3, R["strat"]-0.28, 1.5, 0.36,
             "7k — noisy\nAUC = 0.505",
             C_WARN, fontsize=6.2, text_color="#333")

    draw_box(ax, 10.9, R["strat"]-0.28, 1.5, 0.36,
             "1.5k — clean\nAUC = 0.744",
             C_RESULT, fontsize=6.2)

    arrow(ax, 9.6, R["strat"]-0.05, 8.7, R["strat"]-0.18)
    arrow(ax, 9.6, R["strat"]-0.05, 10.6, R["strat"]-0.18)

    # ── Legend ────────────────────────────────────────────────────────────────
    legend_items = [
        (C_DATA,    "Data / input"),
        (C_SCRIPT,  "Pipeline script"),
        (C_NEG,     "Negative control"),
        (C_RESULT,  "Clean result"),
        (C_WARN,    "Noisy / heterogeneous"),
        (C_NEUTRAL, "Output / table"),
    ]
    for i, (c, lbl) in enumerate(legend_items):
        rx, ry = 0.18, 8.6 - i * 0.32
        patch = FancyBboxPatch((rx, ry - 0.10), 0.22, 0.20,
                               boxstyle="round,pad=0.02",
                               facecolor=c, edgecolor="#333", linewidth=0.6,
                               zorder=5)
        ax.add_patch(patch)
        ax.text(rx + 0.29, ry, lbl, va="center", fontsize=6.2, color="#333")

    ax.set_title("Chameleon Predictor — Pipeline Methodology",
                 fontsize=12, fontweight="bold", pad=6, color="#222")

    plt.tight_layout(pad=0.3)
    plt.savefig(out_path, dpi=180, bbox_inches="tight", facecolor=C_BG)
    plt.close()
    print(f"Flowchart saved: {out_path}")


# ── PowerPoint builder ────────────────────────────────────────────────────────

def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_image_slide(prs, title_text, img_path, caption=None):
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)

    # Title text box
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1),
                                   Inches(9.1), Inches(0.5))
    tf  = txb.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Image — fill most of slide below title
    img_top    = Inches(0.65)
    img_height = Inches(6.2) if caption else Inches(6.6)
    slide.shapes.add_picture(str(img_path),
                             Inches(0.15), img_top,
                             width=Inches(9.4), height=img_height)

    if caption:
        ctxb = slide.shapes.add_textbox(Inches(0.3), Inches(6.95),
                                        Inches(9.1), Inches(0.45))
        ctf  = ctxb.text_frame
        ctf.text = caption
        ctf.paragraphs[0].font.size    = Pt(9)
        ctf.paragraphs[0].font.italic  = True
        ctf.paragraphs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        ctf.paragraphs[0].alignment    = PP_ALIGN.CENTER

    return slide


def build_pptx(outdir: Path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1 — title
    add_title_slide(
        prs,
        "Chameleon Predictor",
        "3D Conformational Descriptors for Cyclic Peptide Membrane Permeation\n"
        "Jorge Carmona · March 2026",
    )

    # Slide 2 — methodology flowchart
    fc_path = outdir / "methodology_flowchart.png"
    if fc_path.exists():
        add_image_slide(prs, "Pipeline Methodology", fc_path)

    # Slide 3 — Panel B 1.5k
    pb_1502 = outdir / "figures" / "Panel_B_3D_delta_umap_1502.png"
    if pb_1502.exists():
        add_image_slide(
            prs, "UMAP Panel B — 3D Δ Features (1,502 compounds, Furukawa + Chugai)",
            pb_1502,
            caption="AUC = 0.744 on clean homogeneous subset. "
                    "Two-population structure: chameleonic (high ΔPSA) vs. rigid/polar.",
        )

    # Slide 4 — Panel B 7k
    pb_7k = outdir / "figures" / "Panel_B_3D_delta_umap.png"
    if pb_7k.exists():
        add_image_slide(
            prs, "UMAP Panel B — 3D Δ Features (Full 7k)",
            pb_7k,
            caption="AUC = 0.505. Two-population structure survives but PAMPA label "
                    "noise (Townsend/Kelly pooled protocol) washes out the permeability signal.",
        )

    # Slide 5 — Panel C 1.5k + Track D
    pc_1566 = outdir / "figures" / "Panel_C_combined_umap_1566.png"
    if pc_1566.exists():
        add_image_slide(
            prs, "UMAP Panel C + Track D — MW Coloring (1,566 compounds)",
            pc_1566,
            caption="Permeable cluster median MW = 1,180 Da vs. impermeable = 820 Da (1.44×). "
                    "Chameleonic behavior is size-gated: large macrolides only.",
        )

    # Slide 6 — Panel C 7k + Track D
    pc_7k = outdir / "figures" / "Panel_C_combined_umap_6938.png"
    if pc_7k.exists():
        add_image_slide(
            prs, "UMAP Panel C + Track D — MW Coloring (Full 7k)",
            pc_7k,
            caption="MW gap disappears on full dataset (permeable = impermeable = 820 Da). "
                    "Cross-source label noise inverts the size-permeability signal.",
        )

    # Slide 7 — AUC bar: Furukawa + Chugai 1,566 (all sizes)
    auc_1566 = outdir / "figures" / "auc_roc_bar_2016_2013_1566.png"
    if auc_1566.exists():
        add_image_slide(
            prs, "AUC-ROC — Furukawa + Chugai (n=1,566, all sizes)",
            auc_1566,
            caption="On clean single-protocol sources, ΔPSA AUC=0.69 and psa3d_std AUC=0.69 — "
                    "matching NumHDonors (0.69) and outperforming MolLogP (0.68). "
                    "3D conformational descriptors are competitive with 2D baselines. "
                    "Rg and mem_hb_count emerge as new predictors for the lipophilic regime.",
        )

    # Slide 8 — AUC bar: Furukawa + Chugai ≥9 residues (chameleonic regime)
    auc_829 = outdir / "figures" / "auc_roc_bar_2016_2013_res9plus_829.png"
    if auc_829.exists():
        add_image_slide(
            prs, "AUC-ROC — ≥9 Residues Only (n=829, chameleonic regime)",
            auc_829,
            caption="Applying the ≥9 residue filter: MolLogP AUC=0.90, MolWt=0.87. "
                    "Rg (membrane conformer) AUC=0.82 — strongest 3D signal yet, above TPSA (0.77). "
                    "ΔPSA falls to 0.44 but norm_delta_psa and NPR1/NPR2 emerge. "
                    "Large chameleonic macrolides: compact 3D shape drives permeability, not switching magnitude alone.",
        )

    # Slide 9 — CREMP overlap AUC bar
    auc_cremp = outdir / "figures" / "auc_roc_bar_cremp_overlap_2416.png"
    if auc_cremp.exists():
        add_image_slide(
            prs, "AUC-ROC — CREMP × CycPeptMPDB Overlap (n=2,416, 6–7 mers)",
            auc_cremp,
            caption="aq_psa3d AUC=0.667 — strongest 3D descriptor, just below MolLogP (0.684). "
                    "bw_psa3d AUC=0.607 (lipophilicity proxy). "
                    "delta_psa3d AUC=0.588 — above chance but weak; 6–7 mers are not switching enough for ΔPSA to dominate. "
                    "pop_lowest_pct AUC=0.520 — conformational rigidity carries no signal in this size regime.",
        )

    # Slide 10 — CREMP structural bias
    cremp_bias = outdir / "figures" / "cremp_structural_bias.png"
    if cremp_bias.exists():
        add_image_slide(
            prs, "CREMP Benchmark — Structural Bias: Wrong Size Regime",
            cremp_bias,
            caption="CREMP overlap (n=2,435) peaks at 6–7 mers — entirely below the ≥9 residue chameleonic threshold (Yu 2026). "
                    "Non-overlap CycPeptMPDB compounds peak at 10 mers, the correct regime. "
                    "ΔPSA cannot be expected to predict permeability on this subset.",
        )

    # Slide 11 — CREMP ROC + two-population summary
    cremp_roc = outdir / "figures" / "cremp_roc_comparison.png"
    if cremp_roc.exists():
        add_image_slide(
            prs, "CREMP Benchmark — ΔPSA at Chance; Lipophilicity Signal Identified",
            cremp_roc,
            caption="Vacuum ΔPSA AUC=0.499, Hybrid ΔPSA (CREMP CHCl₃ mem) AUC=0.454 — both at chance on 6–7 mer overlap. "
                    "bw_psa3d (inverted) AUC=0.607: lower CHCl₃ polar SASA predicts permeability via lipophilicity, not chameleonic switching. "
                    "Two distinct permeability mechanisms confirmed: chameleonic (≥9 mer) vs. lipophilic passive diffusion (<9 mer).",
        )

    out_path = outdir / "Chameleon_Predictor_slides.pptx"
    prs.save(out_path)
    print(f"PowerPoint saved: {out_path}")


# ── Entry point ───────────────────────────────────────────────────────────────

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("--outdir", "-o", default="results")
    return parser.parse_args()


if __name__ == "__main__":
    args   = parse_args()
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    make_flowchart(outdir / "methodology_flowchart.png")
    build_pptx(outdir)
