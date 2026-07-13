# env: chameleon-calc
"""
pipeline_overview_slides.py  (v2 — native PPTX shapes, fully editable)
-----------------------------------------------------------------------
Generates Pipeline_Overview.pptx with 3 slides built entirely from
python-pptx shapes and text boxes — every label is editable in PowerPoint.

Slides:
  1. The Science       — chameleonic ΔPSA, dual-dielectric environments
  2. CREST Pipeline    — 3-step compute flow (crest_v3.1.py)
  3. SLURM Jobs        — 5 parallel jobs via submit_tier2_slurm_updated.py

Usage:
  python scripts/pipeline_overview_slides.py [--outdir .]
"""

import argparse
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
BLUE   = (0x2D, 0x6D, 0xB5)
PURPLE = (0x6B, 0x4F, 0xA0)
GREEN  = (0x27, 0xAE, 0x60)
TEAL   = (0x16, 0xA0, 0x85)
ORANGE = (0xE6, 0x7E, 0x22)
RED    = (0xC0, 0x39, 0x2B)
GREY   = (0x95, 0xA5, 0xA6)
DARK   = (0x2C, 0x3E, 0x50)
WHITE  = (0xFF, 0xFF, 0xFF)
WATER  = (0x21, 0x86, 0xC6)
MEM    = (0xC0, 0x39, 0x2B)
DKTEXT = (0x1A, 0x25, 0x2F)

W, H = 13.33, 7.5   # slide size inches (16:9 widescreen)

ROUNDED_RECT = 5    # msoShapeRoundedRectangle
RECTANGLE    = 1    # msoShapeRectangle


def _rgb(t):
    return RGBColor(*t)


def _box(slide, x, y, w, h, text, fill, fg=WHITE, size=9,
         bold=False, italic=False, shape_id=ROUNDED_RECT, border_color=None):
    """Rounded rectangle with centred, vertically-anchored text."""
    shape = slide.shapes.add_shape(
        shape_id, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    bc = border_color or (0x33, 0x33, 0x33)
    shape.line.color.rgb = _rgb(bc)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(fg)
        run.font.bold = bold and i == 0
        run.font.italic = italic
        p.alignment = PP_ALIGN.CENTER
    return shape


def _txt(slide, x, y, w, h, text, color=DKTEXT, size=9,
         align=PP_ALIGN.CENTER, bold=False, italic=False):
    """Plain floating text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(color)
        run.font.bold = bold and i == 0
        run.font.italic = italic
        p.alignment = align
    return txb


def _arrow(slide, x1, y1, x2, y2, color=GREY, width=1.5):
    """Straight arrow connector (arrowhead at end)."""
    try:
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        ctype = MSO_CONNECTOR_TYPE.STRAIGHT
    except Exception:
        ctype = 1
    conn = slide.shapes.add_connector(
        ctype, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = _rgb(color)
    conn.line.width = Pt(width)
    try:
        ln = conn.line._ln
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'arrow')
        tail.set('w', 'med')
        tail.set('len', 'med')
    except Exception:
        pass
    return conn


def _title(slide, text):
    """Full-width dark title bar."""
    _box(slide, 0, 0, W, 0.62, text, DARK, WHITE,
         size=15, bold=True, shape_id=RECTANGLE, border_color=DARK)


# ── Slide 1: The Science ──────────────────────────────────────────────────────
def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _title(slide, "What is a Chameleonic Cyclic Peptide?")

    # ── Left: problem ─────────────────────────────────────────────────────────
    _txt(slide, 0.2, 0.72, 3.8, 0.35, "The Challenge",
         color=DARK, size=12, bold=True)
    _txt(slide, 0.2, 1.1, 3.8, 1.1,
         "Large cyclic peptides (≥9 residues) must cross\n"
         "cell membranes despite having many polar\n"
         "H-bond donors and acceptors.",
         color=DKTEXT, size=9.5)
    _box(slide, 0.2, 2.25, 3.8, 0.58,
         "\"Beyond Rule of 5\"  —  MW > 500 Da, HBD > 5",
         ORANGE, size=9.5)
    _txt(slide, 0.2, 2.9, 3.8, 0.8,
         "Classical Lipinski rules predict these molecules\n"
         "should NOT be permeable — yet some are.",
         color=GREY, size=9, italic=True)

    # ── Center: switching ─────────────────────────────────────────────────────
    _txt(slide, 4.1, 0.72, 5.2, 0.35, "Conformational Switching",
         color=DARK, size=12, bold=True)

    _box(slide, 4.2, 1.1, 2.3, 0.72,
         "Water  ε = 80\nopen / polar\nconformer", WATER, size=9)
    _box(slide, 7.0, 1.1, 2.3, 0.72,
         "Membrane  ε = 4.8\nclosed / nonpolar\nconformer", MEM, size=9)

    # double-headed arrow
    _txt(slide, 6.35, 1.32, 0.75, 0.3, "⇄", color=PURPLE, size=18)

    _box(slide, 4.2, 2.05, 2.3, 0.55,
         "High PSA  ~180 Å²\n(H-bonds exposed)", WATER, size=9)
    _box(slide, 7.0, 2.05, 2.3, 0.55,
         "Low PSA  ~105 Å²\n(H-bonds buried)", MEM, size=9)

    _arrow(slide, 5.37, 1.82, 5.37, 2.05)
    _arrow(slide, 8.15, 1.82, 8.15, 2.05)

    _box(slide, 5.0, 2.85, 3.5, 0.68,
         "ΔPSA = PSA(water) − PSA(mem)\n= +75 Å²  →  permeable!",
         GREEN, size=10)
    _arrow(slide, 5.37, 2.6, 5.9, 2.85)
    _arrow(slide, 8.15, 2.6, 7.6, 2.85)

    _txt(slide, 4.2, 3.63, 5.2, 0.42,
         "e.g. Cyclosporin A — gold-standard chameleonic peptide",
         color=GREY, size=8.5, italic=True)

    # Mechanism box
    _box(slide, 4.1, 4.15, 5.2, 0.75,
         "How:  intramolecular H-bonds form in membrane\n"
         "→ polar groups buried  →  nonpolar surface exposed",
         PURPLE, size=9)

    # ── Right: reference compounds ────────────────────────────────────────────
    _txt(slide, 9.5, 0.72, 3.65, 0.35, "5 Reference Compounds",
         color=DARK, size=12, bold=True)

    cpds = [
        ("1NMe3  (6-mer)        PAMPA = −5.52",  "permeable — N-methylation, not switching  ✓", TEAL),
        ("CsA  (11-mer)         PAMPA = −5.90",  "permeable — chameleonic  ★",                  GREEN),
        ("c*[PSLYF]  (11-mer)   PAMPA = −9.10",  "impermeable  ✗",                              RED),
        ("DP-955  (15-mer)      PAMPA = −5.20",  "permeable  ✓",                                GREEN),
        ("DP-944  (15-mer)      PAMPA = −7.00",  "impermeable  ✗",                              RED),
    ]
    for i, (name, label, col) in enumerate(cpds):
        _box(slide, 9.5, 1.1 + i * 0.62, 3.65, 0.55,
             f"{name}\n{label}", col, size=8.5)

    _txt(slide, 9.5, 4.25, 3.65, 0.7,
         "1NMe3 = N-methylation control:\npermeable via reduced backbone polarity,\n"
         "not conformational switching  (Bockus, Lokey lab 2015)",
         color=GREY, size=8, italic=True)

    # ── Bottom goal ───────────────────────────────────────────────────────────
    _box(slide, 0.2, 6.7, W - 0.4, 0.65,
         "Goal:  CREST ΔPSA should be HIGH for CsA/DP-955 (chameleonic), "
         "LOW for 1NMe3 (N-methylation), and LOW for impermeable compounds.",
         PURPLE, size=10, bold=True)

    return slide


# ── Slide 2: CREST Pipeline ───────────────────────────────────────────────────
def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Tier-2 CREST Pipeline  —  crest_v3.1.py  (CREST v2.12 + ALPB)")

    # ── Left column: 3-step flow ──────────────────────────────────────────────
    _box(slide, 0.3, 0.75, 3.9, 0.58,
         "Input: SMILES string\n(RDKit computes formal charge automatically)",
         BLUE, size=9)
    _arrow(slide, 2.25, 1.33, 2.25, 1.55)

    _box(slide, 0.3, 1.55, 3.9, 0.75,
         "Step 1 — RDKit ETKDGv3\nEmbed 5,000 conformers → MMFF94 opt\nRMSD filter → top 50",
         BLUE, size=9, bold=True)
    _txt(slide, 0.3, 2.35, 3.9, 0.55,
         "useMacrocycleTorsions=True  |  pruneRmsThresh=0.01",
         color=GREY, size=8, italic=True)
    _arrow(slide, 2.25, 2.3, 2.25, 2.62)

    _box(slide, 0.3, 2.62, 3.9, 0.75,
         "Step 2 — GFN2-xTB --opt\n50 parallel workers, one per conformer\n--alpb {solvent}  OMP_NUM_THREADS=1,1",
         TEAL, size=9, bold=True)
    _txt(slide, 0.3, 3.42, 3.9, 0.4,
         "→ select lowest-energy conformer as CREST seed",
         color=GREY, size=8, italic=True)
    _arrow(slide, 2.25, 3.37, 2.25, 3.68)

    _box(slide, 0.3, 3.68, 3.9, 0.75,
         "Step 3 — CREST iMTD-GC\ncrest {xyz} --gfn2 --alpb {solvent}\n-T {cpus} --keepdir",
         PURPLE, size=9, bold=True)
    _txt(slide, 0.3, 4.48, 3.9, 0.4,
         "Full iMTD-GC, no --quick — matches CREMP protocol exactly",
         color=GREY, size=8, italic=True)
    _arrow(slide, 2.25, 4.43, 2.25, 4.72)

    _box(slide, 0.3, 4.72, 3.9, 0.55,
         "crest_conformers.xyz\n(hundreds of conformers + GFN2 energies in comment)",
         GREY, size=8.5)

    # ── Right: dual-env + ΔPSA ────────────────────────────────────────────────
    _txt(slide, 4.4, 0.72, 8.7, 0.35,
         "Same 3-step pipeline runs independently in each dielectric environment",
         color=DARK, size=9.5, bold=True)

    _box(slide, 4.5, 1.1, 3.6, 0.65,
         "Water ensemble\nε = 80  (--alpb water)", WATER, size=10, bold=True)
    _box(slide, 9.3, 1.1, 3.6, 0.65,
         "Membrane ensemble\nε = 4.8  (--alpb chcl3)", MEM, size=10, bold=True)

    for cx in (6.3, 11.1):
        _arrow(slide, cx, 1.75, cx, 2.05)

    _box(slide, 4.5, 2.05, 3.6, 0.75,
         "Per conformer:\n3D PSA via rdFreeSASA (Bondi radii)\nIntramolecular H-bonds counted",
         WATER, fg=WHITE, size=8.5)
    _box(slide, 9.3, 2.05, 3.6, 0.75,
         "Per conformer:\n3D PSA via rdFreeSASA (Bondi radii)\nIntramolecular H-bonds counted",
         MEM, fg=WHITE, size=8.5)

    for cx in (6.3, 11.1):
        _arrow(slide, cx, 2.8, cx, 3.1)

    _box(slide, 4.5, 3.1, 3.6, 0.78,
         "Boltzmann weighting\nw_i ∝ exp(−E_i / RT)  T = 298.15 K\n→ PSA_boltz,  HB_boltz",
         BLUE, size=8.5)
    _box(slide, 9.3, 3.1, 3.6, 0.78,
         "Boltzmann weighting\nw_i ∝ exp(−E_i / RT)  T = 298.15 K\n→ PSA_boltz,  HB_boltz",
         BLUE, size=8.5)

    _arrow(slide, 6.3, 3.88, 7.5, 4.45)
    _arrow(slide, 11.1, 3.88, 9.9, 4.45)

    _box(slide, 6.0, 4.45, 5.4, 0.75,
         "ΔPSA = PSA_boltz(water) − PSA_boltz(mem)\n"
         "ΔHB  = HB_boltz(mem) − HB_boltz(water)",
         GREEN, size=10, bold=True)

    _arrow(slide, 8.7, 5.2, 8.7, 5.5)

    _box(slide, 6.0, 5.5, 5.4, 0.55,
         "CSV output → results/crest_runs/run_{timestamp}_{idx}_{short}/",
         GREY, size=8.5)

    # version note
    _box(slide, 0.3, 6.8, W - 0.6, 0.48,
         "CREST v2.12  (conda env: chameleon_crest212)  —  downgraded from 3.x "
         "due to crash reproducibility issues on cluster",
         RED, size=9, italic=True)

    return slide


# ── Slide 3: SLURM Parallelization ───────────────────────────────────────────
def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "HPC Parallelization  —  submit_tier2_slurm_updated.py  (SLURM)")

    # ── Submit script ─────────────────────────────────────────────────────────
    _box(slide, 0.3, 0.72, 5.5, 0.65,
         "submit_tier2_slurm_updated.py\npython scripts/submit_tier2_slurm_updated.py [--compounds 0 1 2 3 4]",
         PURPLE, size=9, bold=True)
    _txt(slide, 0.3, 1.42, 5.5, 0.4,
         "--cpus 20   --mem 16G   --partition all   conda env: chameleon_crest212",
         color=GREY, size=8, italic=True)

    # Fan-out arrows to 5 jobs
    job_xs = [1.0, 3.45, 5.9, 8.35, 10.8]
    for jx in job_xs:
        _arrow(slide, 3.05, 1.37, jx + 1.1, 2.05, color=GREY, width=1.0)

    # ── 5 job columns ─────────────────────────────────────────────────────────
    compounds = [
        ("1NMe3\n6-mer",  "permeable  ✓", TEAL),
        ("CsA\n11-mer",       "permeable  ★",  GREEN),
        ("PSLYF\n11-mer",     "impermeable  ✗", RED),
        ("DP-955\n15-mer",    "permeable  ✓",   GREEN),
        ("DP-944\n15-mer",    "impermeable  ✗", RED),
    ]
    for i, (jx, (cname, clabel, ccol)) in enumerate(zip(job_xs, compounds)):
        # Job header
        _box(slide, jx, 2.05, 2.2, 0.65,
             f"Job {i}  —  {cname}\n{clabel}", ccol, size=8.5)

        # SBATCH config
        _arrow(slide, jx + 1.1, 2.7, jx + 1.1, 2.92)
        _box(slide, jx, 2.92, 2.2, 0.62,
             "#SBATCH --cpus=20\n--mem 16G  --ntasks 1",
             DARK, size=8)

        # Script call
        _arrow(slide, jx + 1.1, 3.54, jx + 1.1, 3.75)
        _box(slide, jx, 3.75, 2.2, 0.62,
             f"crest_v3.1.py\n--compound {i} --threads 20",
             BLUE, size=8)

        # Output folder
        _arrow(slide, jx + 1.1, 4.37, jx + 1.1, 4.57)
        _box(slide, jx, 4.57, 2.2, 0.62,
             f"results/crest_runs/\nrun_*_{i}_{cname.split(chr(10))[0]}/",
             GREY, size=7.8)

    # ── Notes ─────────────────────────────────────────────────────────────────
    _box(slide, 0.3, 5.38, W - 0.6, 0.5,
         "All 5 jobs run simultaneously on the cluster  —  "
         "each compound gets its own SLURM allocation",
         GREEN, size=9.5, bold=True)

    _box(slide, 0.3, 5.97, 6.3, 0.58,
         "Logs → results/slurm_logs/run_{timestamp}/\n"
         "crest_{idx}_{short}_{jobid}.out  /  .err",
         TEAL, size=8.5)

    _box(slide, 6.8, 5.97, 6.2, 0.58,
         "Restart mode (--restart): reloads full_ensemble.xyz\n"
         "re-analyze at higher --max-confs without re-running CREST",
         BLUE, size=8.5)

    _txt(slide, 0.3, 6.63, W - 0.6, 0.4,
         "Monitor: squeue -u $USER   |   "
         "Saved ensembles: results/conformers/{short}/{aq|mem}/full_ensemble.xyz",
         color=GREY, size=8, italic=True)

    return slide


# ── Build PPTX ────────────────────────────────────────────────────────────────
def build_pptx(out_path: Path):
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)

    prs.save(out_path)
    print(f"Saved: {out_path}")


def parse_args():
    p = argparse.ArgumentParser(
        description="Generate editable 3-slide pipeline overview PPTX")
    p.add_argument("--outdir", "-o", default=".",
                   help="Output directory (default: current directory)")
    return p.parse_args()


if __name__ == "__main__":
    args    = parse_args()
    outdir  = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    build_pptx(outdir / "Pipeline_Overview.pptx")
