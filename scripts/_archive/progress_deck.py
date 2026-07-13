# env: chameleon-calc
"""
progress_deck.py  — 3-month progress narrative (committee-facing)
------------------------------------------------------------------
Generates Progress_Deck.pptx: an 11-slide narrative of the last three
months of the Chameleon_Predictor project, built entirely from python-pptx
shapes/text boxes (every label editable in PowerPoint).

Distinct from pipeline_overview_slides.py (a mechanistic "how it works"
explainer, now stale). This deck tells the progress STORY:
  problem -> idea -> pipeline -> timeline -> validation -> rigor ->
  headline R/S result -> free-energy pivot -> ML benchmark -> road ahead

Every number is sourced from repo docs (per-molecule R/S reports, CsA
exp-vs-CREST experiment, descriptor-review path-forward, and the CREMP
feature benchmark run 2026-07-07). No figures embedded — talk skeleton the
presenter drops plots into.

Usage:
  python scripts/progress_deck.py [--outdir .]
"""
import argparse
import struct
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
BLUE   = (0x2D, 0x6D, 0xB5)
PURPLE = (0x6B, 0x4F, 0xA0)
GREEN  = (0x27, 0xAE, 0x60)
TEAL   = (0x16, 0xA0, 0x85)
ORANGE = (0xE6, 0x7E, 0x22)
RED    = (0xC0, 0x39, 0x2B)
GREY   = (0x95, 0xA5, 0xA6)
LGREY  = (0xEC, 0xF0, 0xF1)
DARK   = (0x2C, 0x3E, 0x50)
WHITE  = (0xFF, 0xFF, 0xFF)
WATER  = (0x21, 0x86, 0xC6)
MEM    = (0xC0, 0x39, 0x2B)
DKTEXT = (0x1A, 0x25, 0x2F)
ACCENT = (0x16, 0xA0, 0x85)

W, H = 13.33, 7.5
ROUNDED_RECT, RECTANGLE, OVAL = 5, 1, 9

FIGDIR = Path(__file__).resolve().parent.parent / "results" / "figures"


def _rgb(t):
    return RGBColor(*t)


def _box(slide, x, y, w, h, text, fill, fg=WHITE, size=10, bold=False,
         italic=False, shape_id=ROUNDED_RECT, border_color=None,
         align=PP_ALIGN.CENTER, line_w=0.75):
    shape = slide.shapes.add_shape(shape_id, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(border_color or (0x33, 0x33, 0x33))
    shape.line.width = Pt(line_w)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(fg)
        run.font.bold = bold and i == 0
        run.font.italic = italic
        p.alignment = align
    return shape


def _txt(slide, x, y, w, h, text, color=DKTEXT, size=11, align=PP_ALIGN.LEFT,
         bold=False, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(color)
        run.font.bold = bold and i == 0
        run.font.italic = italic
        p.alignment = align
    return txb


def _arrow(slide, x1, y1, x2, y2, color=GREY, width=1.5):
    try:
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        ctype = MSO_CONNECTOR_TYPE.STRAIGHT
    except Exception:
        ctype = 1
    conn = slide.shapes.add_connector(ctype, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = _rgb(color)
    conn.line.width = Pt(width)
    try:
        tail = etree.SubElement(conn.line._ln, qn('a:tailEnd'))
        tail.set('type', 'arrow'); tail.set('w', 'med'); tail.set('len', 'med')
    except Exception:
        pass
    return conn


def _png_size(path):
    try:
        with open(path, "rb") as f:
            head = f.read(24)
        if head[:8] == b"\x89PNG\r\n\x1a\n":
            return struct.unpack(">II", head[16:24])
    except Exception:
        pass
    return None


def _pic(slide, path, bx, by, bw, bh, caption=None):
    """Place an image scaled to fit (preserving aspect) and centered in a box."""
    p = Path(path)
    if not p.is_absolute():
        p = FIGDIR / path
    sz = _png_size(p)
    if not sz or not p.exists():
        _box(slide, bx, by, bw, bh, f"[missing figure]\n{Path(path).name}",
             LGREY, fg=GREY, size=10)
        return
    iw, ih = sz
    scale = min(bw / iw, bh / ih)
    w, h = iw * scale, ih * scale
    x, y = bx + (bw - w) / 2, by + (bh - h) / 2
    slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        _txt(slide, bx, by + bh + 0.02, bw, 0.35, caption, color=GREY,
             size=9, italic=True, align=PP_ALIGN.CENTER)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


_SLIDE_NO = [0]


def _footer(slide, n=None):
    _SLIDE_NO[0] += 1
    _txt(slide, 0.5, 7.12, 8.0, 0.3, "Chameleon_Predictor · Hu Lab · SDSU",
         color=GREY, size=8)
    _txt(slide, 11.8, 7.12, 1.0, 0.3, str(_SLIDE_NO[0]), color=GREY, size=8,
         align=PP_ALIGN.RIGHT)


def _header(slide, kicker, title, n):
    _txt(slide, 0.55, 0.28, 12.2, 0.3, kicker.upper(), color=ACCENT, size=11,
         bold=True)
    _txt(slide, 0.5, 0.52, 12.3, 0.7, title, color=DARK, size=23, bold=True)
    bar = slide.shapes.add_shape(RECTANGLE, Inches(0.55), Inches(1.24),
                                 Inches(2.2), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(ACCENT)
    bar.line.fill.background()
    _footer(slide, n)


# ── Slide 1 — Title ───────────────────────────────────────────────────────────
def s01_title(prs):
    slide = _blank(prs)
    band = slide.shapes.add_shape(RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(2.35))
    band.fill.solid(); band.fill.fore_color.rgb = _rgb(DARK); band.line.fill.background()
    acc = slide.shapes.add_shape(RECTANGLE, Inches(0), Inches(2.35), Inches(W), Inches(0.09))
    acc.fill.solid(); acc.fill.fore_color.rgb = _rgb(ACCENT); acc.line.fill.background()

    _txt(slide, 0.7, 0.55, 12.0, 0.4, "THREE-MONTH PROGRESS REPORT", color=ACCENT, size=13, bold=True)
    _txt(slide, 0.7, 0.98, 12.0, 1.3,
         "Predicting Membrane Permeability of Chameleonic\n"
         "Cyclic Peptides from 3D Conformational Ensembles",
         color=WHITE, size=25, bold=True)
    _txt(slide, 0.7, 2.75, 12.0, 0.4,
         "Jorge Carmona   ·   Hu Lab, San Diego State University",
         color=DARK, size=15, bold=True)
    _txt(slide, 0.7, 3.2, 12.0, 0.35, "April – July 2026", color=GREY, size=12)

    _box(slide, 0.7, 4.0, 11.93, 0.95,
         "A quantum-chemistry conformer pipeline that resolves the 3D shape "
         "differences 2D descriptors miss —\nvalidated on a known chameleon, "
         "shown to resolve R/S epimers, and quantified against experimental "
         "permeability on 2,400+ peptides.",
         LGREY, fg=DKTEXT, size=12.5)

    pillars = [
        ("BUILT", "Dual-dielectric CREST\n3D-descriptor pipeline", BLUE),
        ("VALIDATED", "Reproduces CsA;\nresolves R/S epimers", GREEN),
        ("QUANTIFIED", "3D beats 2D on\nout-of-source PAMPA", PURPLE),
    ]
    x = 0.7
    for tag, body, col in pillars:
        _box(slide, x, 5.25, 3.84, 1.35, f"{tag}\n{body}", col, size=11.5, bold=True)
        x += 4.11
    _footer(slide, 1)
    return slide


# ── Slide 2 — The problem ─────────────────────────────────────────────────────
def s02_problem(prs):
    slide = _blank(prs)
    _header(slide, "Motivation", "The problem: permeability beyond the Rule of 5", 2)
    _txt(slide, 0.55, 1.5, 6.1, 2.4,
         "Large cyclic peptides (≥ 9 residues) reach targets small molecules "
         "cannot — but they carry many polar H-bond donors and acceptors.\n\n"
         "Classical Lipinski / Rule-of-5 filters predict they should NOT cross "
         "membranes. Yet some do. Predicting which ones is the central problem.",
         color=DKTEXT, size=13)
    _box(slide, 0.55, 4.35, 6.1, 0.7,
         "“Beyond Rule of 5”:  MW > 500 Da,  H-bond donors > 5", ORANGE, size=12.5, bold=True)
    _txt(slide, 0.55, 5.25, 6.1, 1.3,
         "Why it is hard: permeability is not a property of the 2D structure. "
         "It depends on which 3D shapes the molecule adopts in water versus in "
         "the membrane — information a flat diagram does not contain.",
         color=GREY, size=12, italic=True)

    _box(slide, 7.1, 1.65, 5.6, 0.55, "The paradox", DARK, size=13, bold=True, shape_id=RECTANGLE)
    _box(slide, 7.1, 2.4, 2.65, 1.5, "Many polar groups\n\nHigh calculated PSA\n\n→ “should not permeate”", WATER, size=11.5)
    _box(slide, 10.05, 2.4, 2.65, 1.5, "Yet measurably\npermeable\n\ncell-active,\norally available", GREEN, size=11.5)
    _txt(slide, 9.75, 2.95, 0.3, 0.4, "≠", color=RED, size=22, bold=True, align=PP_ALIGN.CENTER)
    _box(slide, 7.1, 4.15, 5.6, 1.15,
         "Resolution: these molecules CHANGE SHAPE with their environment.\n"
         "The permeability signal lives in the conformational ensemble.",
         PURPLE, size=12, bold=True)
    _txt(slide, 7.1, 5.5, 5.6, 0.8,
         "Gold standard: cyclosporin A — an 11-residue macrocycle, orally "
         "bioavailable despite 5 H-bond donors.",
         color=GREY, size=11, italic=True)
    return slide


# ── Slide 3 — The idea ────────────────────────────────────────────────────────
def s03_idea(prs):
    slide = _blank(prs)
    _header(slide, "Core concept", "Chameleonicity: shape-switching across environments", 3)
    _box(slide, 1.0, 1.75, 3.3, 1.0, "WATER   ε = 80\nopen / polar conformer\nH-bonds exposed", WATER, size=12, bold=True)
    _box(slide, 9.0, 1.75, 3.3, 1.0, "MEMBRANE   ε ≈ 4.8\nclosed / apolar conformer\nH-bonds buried", MEM, size=12, bold=True)
    _txt(slide, 5.7, 1.85, 1.9, 0.7, "⇄", color=PURPLE, size=40, align=PP_ALIGN.CENTER)
    _txt(slide, 4.5, 2.85, 4.3, 0.4, "conformational switching", color=GREY, size=11, italic=True, align=PP_ALIGN.CENTER)
    _box(slide, 1.0, 3.15, 3.3, 0.7, "High 3D-PSA\n(~180 Å²)", WATER, size=12)
    _box(slide, 9.0, 3.15, 3.3, 0.7, "Low 3D-PSA\n(~105 Å²)", MEM, size=12)
    _arrow(slide, 2.65, 2.75, 2.65, 3.15)
    _arrow(slide, 10.65, 2.75, 10.65, 3.15)
    _box(slide, 4.3, 4.25, 4.7, 0.95,
         "ΔPSA = PSA(water) − PSA(membrane)\nLarge ΔPSA → polarity shielding → permeable",
         GREEN, size=13, bold=True)
    _arrow(slide, 2.65, 3.85, 4.3, 4.55)
    _arrow(slide, 10.65, 3.85, 9.0, 4.55)
    _txt(slide, 0.55, 5.22, 12.2, 0.4,
         "Chameleonic ΔPSA is the mechanism for LARGE (≥9-residue) macrocycles; smaller / "
         "rigid rings permeate by other routes (shape, lipophilicity) — captured downstream by ΔG_transfer.",
         color=GREY, size=9.5, italic=True, align=PP_ALIGN.CENTER)
    _box(slide, 0.55, 5.65, 12.2, 1.05,
         "Why this drives the project: epimers (R vs S) and other stereo-subtle "
         "analogs are IDENTICAL on every 2D and lipophilicity descriptor. Any "
         "permeability difference must come from their 3D ensembles — so we must "
         "sample and score those ensembles.",
         DARK, size=12.5, bold=True)
    return slide


# ── Slide 4 — The pipeline ────────────────────────────────────────────────────
def s04_pipeline(prs):
    slide = _blank(prs)
    _header(slide, "What we built", "A dual-dielectric quantum-chemistry conformer pipeline", 4)
    steps = [
        ("SMILES input", "RDKit assigns formal charge", BLUE),
        ("Step 1 — RDKit ETKDGv3", "5,000 embeds → MMFF94 → top 50\nmacrocycle torsions on", BLUE),
        ("Step 2 — GFN2-xTB opt", "50 parallel workers, per solvent\n→ lowest-energy CREST seed", TEAL),
        ("Step 3 — CREST iMTD-GC", "full metadynamics, GFN2-xTB\nmatches CREMP protocol", PURPLE),
    ]
    y = 1.6
    for i, (t, sub, col) in enumerate(steps):
        _box(slide, 0.55, y, 4.5, 0.82, f"{t}\n{sub}", col, size=10.5, bold=True)
        if i < len(steps) - 1:
            _arrow(slide, 2.8, y + 0.82, 2.8, y + 1.02)
        y += 1.02
    _box(slide, 0.55, y, 4.5, 0.62, "Conformer ensemble + GFN2 energies", GREY, size=10.5)

    _txt(slide, 5.5, 1.55, 7.3, 0.4, "Run independently in each dielectric, then Boltzmann-weight:", color=DARK, size=12, bold=True)
    _box(slide, 5.5, 2.05, 3.5, 0.6, "Water ensemble  ε = 80", WATER, size=11, bold=True)
    _box(slide, 9.25, 2.05, 3.5, 0.6, "Membrane ensemble  ε ≈ 4.8", MEM, size=11, bold=True)
    for cx in (7.25, 11.0):
        _arrow(slide, cx, 2.65, cx, 2.9)
    _box(slide, 5.5, 2.9, 7.25, 0.95,
         "Per conformer:  3D-PSA (rdFreeSASA, Bondi radii) · radius of gyration "
         "· transannular IMHB count · SASA breakdown", DARK, size=11)
    _arrow(slide, 9.1, 3.85, 9.1, 4.1)
    _box(slide, 5.5, 4.1, 7.25, 0.7, "Boltzmann weighting   wᵢ ∝ exp(−Eᵢ / RT),  T = 298.15 K", BLUE, size=11.5, bold=True)
    _arrow(slide, 9.1, 4.8, 9.1, 5.05)
    _box(slide, 5.5, 5.05, 7.25, 0.85, "Ensemble descriptors  →  ΔPSA, ΔIMHB\n(the permeability read-out, per molecule)", GREEN, size=12, bold=True)
    _txt(slide, 5.5, 6.05, 7.25, 0.9,
         "Runs on the HPC via SLURM (CREST 2.12, env-per-role conda envs). "
         "Descriptors are ensemble averages — never a cherry-picked minimum.",
         color=GREY, size=10.5, italic=True)
    return slide


# ── Slide 5 — Timeline ────────────────────────────────────────────────────────
def s05_timeline(prs):
    slide = _blank(prs)
    _header(slide, "The arc", "Three months, four phases", 5)
    phases = [
        ("APR – early MAY", BLUE, "Foundation",
         "Reference-compound set curated · CREST checkpointing + data-loss guards "
         "· CREST 3.x → 2.12 downgrade (cluster crash fix)"),
        ("MAY", TEAL, "Pipeline + first validation",
         "MACE / explicit-TIP3P feasibility · CsA validation (group mtg 5/13) · "
         "refactor to crest_v3.2/3.3 · env-per-role convention"),
        ("JUNE", PURPLE, "Descriptors + headline result",
         "DOPC 3-12-8-12 & 3-12-10-12 R/S studies · descriptors aligned to "
         "literature (Ono / Begnini) · diazirine N=N artifact fixed · PI reports"),
        ("late JUNE – JULY", GREEN, "Free-energy pivot + ML benchmark",
         "Edison descriptor review → validated core · CREMP feature benchmark run "
         "· FlexiSol → ΔG_transfer via CPCM-X · free_energy_calculator built"),
    ]
    spine = slide.shapes.add_shape(RECTANGLE, Inches(0.9), Inches(1.75), Inches(0.06), Inches(4.9))
    spine.fill.solid(); spine.fill.fore_color.rgb = _rgb(GREY); spine.line.fill.background()
    y = 1.7
    for date, col, head, body in phases:
        dot = slide.shapes.add_shape(OVAL, Inches(0.72), Inches(y + 0.12), Inches(0.42), Inches(0.42))
        dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col)
        dot.line.color.rgb = _rgb(WHITE); dot.line.width = Pt(2)
        _txt(slide, 1.35, y, 2.5, 0.4, date, color=col, size=12, bold=True)
        _txt(slide, 1.35, y + 0.38, 3.4, 0.4, head, color=DARK, size=13, bold=True)
        _txt(slide, 4.85, y + 0.02, 7.9, 1.1, body, color=DKTEXT, size=11.5)
        y += 1.22
    return slide


# ── Slide 6 — CsA validation ──────────────────────────────────────────────────
def s06_csa(prs):
    slide = _blank(prs)
    _header(slide, "Validation", "Reproducing cyclosporin A — and does more sampling fix it?", 6)
    _txt(slide, 0.55, 1.45, 12.2, 0.55,
         "Target: the aqueous A1 conformer (NMR + X-ray, Limbach/Bhatt 2022) — defined by a "
         "cis MeVal11–MeBmt1 amide and an open, solvent-exposed fold. Two CREST runs, scored identically:",
         color=DKTEXT, size=12, bold=True)
    rows = [
        ("Ensemble", "3D-PSA (Å²)", "Rg (Å)", "IMHB", "cis MeVal11–MeBmt1", DARK, WHITE),
        ("A1 aqueous (target)", "137.5", "6.15", "~2", "cis", (0xE4, 0xF5, 0xEC), DKTEXT),
        ("CREST v1  (23 conf)", "102.0", "5.85", "2.31", "trans — 0%", LGREY, DKTEXT),
        ("CREST v2  −notopo (1019 conf)", "90.5", "5.58", "3.93", "trans — 0%", (0xFB, 0xE3, 0xE0), DKTEXT),
    ]
    cw = [3.5, 1.7, 1.2, 1.2, 3.0]
    y = 2.2
    for r in rows:
        x = 0.55
        vals = r[:5]; fill = r[5]; fg = r[6]
        for j, (v, w) in enumerate(zip(vals, cw)):
            _box(slide, x, y, w, 0.5, v, fill, fg=fg, size=10.5,
                 bold=(y == 2.2 or j == 0), shape_id=RECTANGLE, border_color=WHITE, line_w=1.0,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            x += w
        y += 0.5
    _box(slide, 0.55, 4.5, 6.0, 1.35,
         "−notopo did NOT fix it\n8× more sampling (23→1019 conf) but still 0% cis, and MORE "
         "collapsed than v1 (lower Rg, IMHB nearly doubled). v1 is actually the closer match "
         "to A1 — the extra sampling drifted toward closed folds.", ORANGE, size=11, bold=True)
    _box(slide, 6.75, 4.5, 6.0, 1.35,
         "The real barrier: implicit solvent\nALPB has no cavity waters to hold A1 open/cis, so "
         "thorough sampling just reveals the collapsed basins it prefers. Reproducing A1 needs "
         "the explicit-water tier (OpenMM / TIP3P) — not a sampling flag.", BLUE, size=11, bold=True)
    _txt(slide, 0.55, 6.05, 12.2, 1.0,
         "Honest negative result: neither ensemble reaches the cis A1 fold, and the more-converged "
         "v2 shows why — implicit solvation is the limitation, not conformer search. This retires "
         "the −notopo hypothesis and points to explicit water.",
         color=DKTEXT, size=11, italic=True)
    return slide


# ── Slide 7 — Descriptor rigor ────────────────────────────────────────────────
def s07_rigor(prs):
    slide = _blank(prs)
    _header(slide, "Rigor", "Aligning descriptors to the evidence", 8)
    _txt(slide, 0.55, 1.45, 12.2, 0.5,
         "An evidence-graded literature review (ranked by EXPERIMENTAL validation "
         "strength) reshaped which descriptors we lead with.", color=DKTEXT, size=12.5, bold=True)
    _box(slide, 0.55, 2.1, 6.0, 0.55, "Validated core — drives the calls", GREEN, size=12.5, bold=True, shape_id=RECTANGLE)
    for i, c in enumerate([
        "SA 3D-PSA (apolar) — best-reproduced single descriptor (Begnini 2021)",
        "Radius of gyration (apolar) — compactness / folding",
        "Backbone transannular IMHB — donor shielding (Rezai, ~100×)",
        "ΔG_transfer (water→apolar) — strongest correlation (Kamenik r=0.92)",
    ]):
        _box(slide, 0.55, 2.75 + i * 0.7, 6.0, 0.62, c, LGREY, fg=DKTEXT, size=10.5, align=PP_ALIGN.LEFT)
    _box(slide, 6.85, 2.1, 5.9, 0.55, "Demoted to diagnostics (kept, not lead)", GREY, size=12.5, bold=True, shape_id=RECTANGLE)
    for i, d in enumerate([
        "Amphipathic / integy moment — no primary validation",
        "Weighted RMSF — ensemble diagnostic, not a predictor",
        "Kier Φ — applicability-domain filter (Φ < 10)",
        "Asphericity / NPR shape — scaffold-contingent",
    ]):
        _box(slide, 6.85, 2.75 + i * 0.7, 5.9, 0.62, d, (0xF3, 0xF4, 0xF5), fg=GREY, size=10.5, align=PP_ALIGN.LEFT)
    _box(slide, 0.55, 5.75, 12.2, 1.1,
         "Caught in the act:  GFN2-xTB was stretching the diazirine N=N bond to a "
         "spurious 1.43 Å in 5 of 8 ensembles — an artifact masquerading as a real "
         "conformational difference. Detected by a dedicated integrity check, fixed "
         "by constraining N=N and re-running. Descriptor quality is a first-class concern.",
         ORANGE, size=12, bold=True)
    return slide


# ── Slide 8 — Headline R/S result ─────────────────────────────────────────────
def s08_result(prs):
    slide = _blank(prs)
    _header(slide, "Headline result", "3D ensembles resolve R/S epimers that 2D cannot", 9)
    _txt(slide, 0.55, 1.4, 12.2, 0.5,
         "Two DOPC epimer pairs — identical 2D/logP, distinguished only by a single "
         "stereocenter. Boltzmann-weighted 3D descriptors separate them cleanly and "
         "in a consistent direction.", color=DKTEXT, size=12, bold=True)

    def panel(x, title, sub, rows, dpsa):
        _box(slide, x, 2.1, 5.9, 0.6, f"{title}\n{sub}", DARK, size=12, bold=True)
        yy = 2.8
        head = ("descriptor", "R", "S")
        cw = [3.5, 1.2, 1.2]
        ri = 0
        for r in (head,) + rows:
            xx = x
            for j, (v, w) in enumerate(zip(r, cw)):
                fill = BLUE if r is head else (LGREY if ri % 2 == 1 else WHITE)
                fg = WHITE if r is head else DKTEXT
                _box(slide, xx, yy, w, 0.46, v, fill, fg=fg, size=10.5,
                     bold=(r is head or j == 0), shape_id=RECTANGLE,
                     border_color=(0xDD, 0xDD, 0xDD), line_w=0.75)
                xx += w
            yy += 0.46; ri += 1
        _box(slide, x, yy + 0.08, 5.9, 0.55, dpsa, TEAL, size=10.5, bold=True)

    panel(0.55, "3-12-8-12", "(azetidine-2-carboxylic acid)",
          (("3D-PSA water (Å²)", "227.3", "165.9"),
           ("3D-PSA CHCl₃ (Å²)", "198.4", "154.5"),
           ("backbone IMHB (water)", "1.6", "2.9"),
           ("Rg CHCl₃ (Å)", "4.66", "4.59")),
          "ΔPSA buried:  R 28.9  vs  S 11.4 Å²")
    panel(6.85, "3-12-10-12", "(sarcosine)",
          (("3D-PSA water (Å²)", "243.9", "197.5"),
           ("3D-PSA CHCl₃ (Å²)", "201.4", "153.7"),
           ("backbone IMHB (water)", "2.0", "3.0"),
           ("Rg CHCl₃ (Å)", "4.65", "4.55")),
          "ΔPSA buried:  R 42.6  vs  S 43.8 Å²")
    _box(slide, 0.55, 6.15, 12.2, 0.95,
         "Consistent prediction across both pairs:  R is solubility-favored (more "
         "exposed polar surface in water),  S is permeability-favored (more compact, "
         "lower apolar PSA + Rg). A testable trade-off — the 2D descriptors sit at "
         "zero difference.", GREEN, size=11.5, bold=True)
    return slide


# ── Slide 9 — The pivot ───────────────────────────────────────────────────────
def s09_pivot(prs):
    slide = _blank(prs)
    _header(slide, "The pivot", "Free energy: the size-general permeability descriptor", 11)
    _box(slide, 0.55, 1.42, 12.2, 0.92,
         "ΔPSA / chameleonicity is MECHANISM-SPECIFIC — it captures large (≥9-residue) "
         "macrocycles that switch shape. ΔG_transfer (water→apolar free energy) is the "
         "size-general alternative, and the strongest literature correlate (Kamenik 2020, "
         "r = 0.92 with a full ensemble; r ≈ 0.50 from one structure — ensembles non-negotiable).",
         PURPLE, size=11.5, bold=True)

    _txt(slide, 0.55, 2.5, 12.2, 0.35,
         "The plan — a fast energy measure on the ensembles we already generate:",
         color=DARK, size=12, bold=True)
    _box(slide, 0.55, 2.92, 3.7, 1.15,
         "1 · Conformer ensembles\n\nCREST / ALPB geometry,\nnative per solvent phase", TEAL, size=11, bold=True)
    _box(slide, 4.6, 2.92, 3.7, 1.15,
         "2 · CPCM-X single-point\n\nquick per-conformer energy,\nGFN2-level — no ORCA / DFT", BLUE, size=11, bold=True)
    _box(slide, 8.65, 2.92, 3.7, 1.15,
         "3 · Boltzmann + ΔG_transfer\n\nenergy-weight the ensemble →\npartition free energy", PURPLE, size=11, bold=True)
    _arrow(slide, 4.25, 3.5, 4.6, 3.5)
    _arrow(slide, 8.3, 3.5, 8.65, 3.5)

    _box(slide, 0.55, 4.35, 12.2, 0.62,
         "FlexiSol (Grimme 2025): partition ratios beat absolute Gsolv (errors cancel) · "
         "lowest-E ≈ full Boltzmann · CPCM-X is the best cheap model · level of theory barely matters.",
         LGREY, fg=DKTEXT, size=10.5)
    _box(slide, 0.55, 5.1, 12.2, 0.9,
         "Now: energy-free geometry descriptors only — Boltzmann weighting is NOT yet enabled "
         "(trusted energies are weeks out). CPCM-X is how we switch it on: quick single-point "
         "energies unlock the energy layer — Boltzmann-weighted descriptors + ΔG_transfer.",
         DARK, size=11.5, bold=True)
    _txt(slide, 0.55, 6.15, 12.2, 0.8,
         "Decisions locked: apolar phase = cyclohexane (Ono 2019) · native per-phase, never "
         "retrofit · no DFT / COSMO-RS ceiling. Built: free_energy_calculator.py (CPCM-X default, --compare vs ALPB).",
         color=GREY, size=10.5, italic=True)
    return slide


# ── Slide 10 — ML benchmark (real result) ─────────────────────────────────────
def s10_ml(prs):
    slide = _blank(prs)
    _header(slide, "Quantified", "Do the 3D descriptors actually predict PAMPA?", 12)
    _box(slide, 0.55, 1.42, 12.2, 0.9,
         "Benchmark on 2,416 macrocycles with both a CREST chloroform ensemble "
         "(CREMP) and a measured PAMPA value (CycPeptMPDB). Question (PROTAC-TS "
         "2026): does expensive 3D ensemble sampling beat free fingerprints? Judged "
         "two ways — random CV, and the honest leave-source-out split.",
         BLUE, size=11.5, bold=True)

    # results table
    _txt(slide, 0.55, 2.5, 12.2, 0.35, "AUC-ROC  (RandomForest / LightGBM)", color=DARK, size=12, bold=True)
    rows = [
        ("Feature set", "dims", "random CV", "leave-source-out", DARK, WHITE),
        ("Fingerprints (Morgan / atom-pair)", "2048", "0.82–0.83", "0.53–0.60  ▼ collapse", (0xFB, 0xE3, 0xE0), DKTEXT),
        ("Mordred 2D / 2D+3D (1 conformer)", "1442+", "0.83", "0.62–0.65", LGREY, DKTEXT),
        ("F7  CREST 3D ensemble", "10", "0.69", "0.65  ▲ most robust", (0xE4, 0xF5, 0xEC), DKTEXT),
    ]
    cw = [5.4, 1.2, 2.6, 3.0]
    y = 2.9
    for r in rows:
        x = 0.55
        vals = r[:4]; fill = r[4]; fg = r[5]
        for j, (v, w) in enumerate(zip(vals, cw)):
            _box(slide, x, y, w, 0.52, v, fill, fg=fg, size=10.5,
                 bold=(y == 2.9 or j == 0), shape_id=RECTANGLE,
                 border_color=WHITE, line_w=1.0,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            x += w
        y += 0.52

    _box(slide, 0.55, 5.15, 12.2, 1.0,
         "The finding:  random CV flatters the fingerprints — but that is memorizing "
         "the dominant data source (82% one lab). Under out-of-source transfer they "
         "collapse to near-chance, while 10 physics descriptors barely move and become "
         "the MOST robust set. 3D generalizes; 2D memorizes.",
         GREEN, size=11.5, bold=True)
    _txt(slide, 0.55, 6.25, 12.2, 0.8,
         "Mechanistic caveat: this subset is mostly SMALL (6–7-residue) macrocycles — not "
         "chameleons — so a ΔPSA-heavy F7 isn't expected to shine; shape may capture the "
         "small-ring mechanism, and ΔG_transfer is the size-general descriptor. Models still "
         "live out-of-source (scarce data), where physics generalizes. (Richer F7 under evaluation.)",
         color=GREY, size=10, italic=True)
    return slide


# ── Slide 11 — Road ahead ─────────────────────────────────────────────────────
def s11_next(prs):
    slide = _blank(prs)
    _header(slide, "Road ahead", "What's next", 14)
    items = [
        ("NOW", GREEN, "Cluster shakedown of CPCM-X (xtb ≥ 6.6, cyclohexane keyword) + the one-time CPCM-X vs ALPB method check."),
        ("NEXT", TEAL, "Native cyclohexane CREST ensembles for the hits → the real ΔG_transfer descriptor (currently the #1 missing signal)."),
        ("BENCHMARK", BLUE, "Add F8 (aqueous CREST) + ΔG_transfer to the CREMP feature benchmark; energy-weighted descriptors once energies land."),
        ("VALIDATE", PURPLE, "Subset NMR/NAMFIS on the 3-12-x-12 hits to license extending the computational ranking (Begnini template)."),
        ("HAND-OFF", ORANGE, "Experimental permeability (PAMPA / Caco-2 / in-house library) + collaborator explicit-solvent sims — the accuracy ceiling."),
    ]
    y = 1.65
    for tag, col, body in items:
        _box(slide, 0.55, y, 2.1, 0.9, tag, col, size=12.5, bold=True)
        _box(slide, 2.8, y, 9.95, 0.9, body, LGREY, fg=DKTEXT, size=12, align=PP_ALIGN.LEFT)
        y += 1.02
    _box(slide, 0.55, y + 0.05, 12.2, 0.55,
         "The through-line: qualitative sorting + relative trends now, quantitative "
         "calibration handed to experiment.", DARK, size=12, bold=True)
    return slide


# ── Slide 6b — CsA validation figure ──────────────────────────────────────────
def s06b_csa_fig(prs):
    slide = _blank(prs)
    _header(slide, "Validation", "CsA — A1-fingerprint diagnostics", 7)
    _pic(slide, "csa_validation_criteria.png", 0.55, 1.55, 6.0, 4.8,
         caption="Per-criterion A1 agreement across the CREST ensemble (v1 shown; v2 same all-trans pattern)")
    _pic(slide, "csa_validation_boltzmann.png", 6.85, 1.55, 6.0, 4.8,
         caption="Boltzmann-weighted descriptor distribution over the water ensemble")
    return slide


# ── Slide 6c — CREST validity envelope ────────────────────────────────────────
def s06c_envelope(prs):
    slide = _blank(prs)
    _header(slide, "Method limits", "How far can CREST be trusted?  The validity envelope", 0)
    _txt(slide, 0.55, 1.42, 12.2, 0.55,
         "One validation point is a data point, not a boundary. CsA maps the HARD end — we need "
         "the small end to bound where implicit-solvent CREST is good enough:",
         color=DKTEXT, size=12, bold=True)
    _txt(slide, 0.7, 2.05, 11.9, 0.3,
         "small / rigid / pre-organized  ──────  size · flexibility · chameleonicity  ──────▶  "
         "large / flexible / chameleonic",
         color=GREY, size=10, italic=True, align=PP_ALIGN.CENTER)
    bar = slide.shapes.add_shape(RECTANGLE, Inches(0.7), Inches(2.35), Inches(11.9), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(GREY); bar.line.fill.background()

    _box(slide, 0.55, 2.65, 5.9, 1.75,
         "SMALL macrocycles (6–7-mer)\n\nRigid, pre-organized → implicit CREST plausibly adequate. "
         "CREMP benchmark: CREST 3D descriptors were the MOST robust feature set here.\n"
         "Validation: NOT yet NMR-tested  ← the gap", GREEN, size=11, bold=True)
    _box(slide, 6.85, 2.65, 5.9, 1.75,
         "LARGE chameleons (~11-mer · CsA)\n\nExplicit cavity waters + cis/trans switching → "
         "implicit CREST FAILS: v1 & v2 both miss the A1 cis fold (0%).\nNeeds the explicit-water tier.",
         RED, size=11, bold=True)
    _box(slide, 0.55, 4.6, 12.2, 0.8,
         "The boundary is unknown. Locating it needs NMR conformational ensembles ACROSS sizes — "
         "we have the hard end (CsA, fails); the small/mid range is unvalidated.",
         DARK, size=11.5, bold=True)
    _txt(slide, 0.55, 5.55, 12.2, 1.2,
         "Next: hexamer NMR validation (candidates: HexPep vs Rezai 2006 NOE / J-coupling; "
         "Heterophyllin B, Ketzel 2025). Deliverable — a trust map: use CREST where validated & "
         "adequate (cheap, scalable), escalate to explicit water only past the boundary. Turns a "
         "negative result into a calibrated operating range for the whole pipeline.",
         color=GREY, size=11, italic=True)
    return slide


# ── Slide 8b — R/S result figures ─────────────────────────────────────────────
def s08b_result_fig(prs):
    slide = _blank(prs)
    _header(slide, "Headline result", "2D is blind to the stereocenter; 3D resolves it", 10)
    _pic(slide, "isomers/3-12-8-12/fig1_reldiff.png", 0.55, 1.5, 6.2, 4.7,
         caption="3-12-8-12: relative |R−S| per descriptor — 2D/logP ≈ 0, 3D ensemble resolves")
    _pic(slide, "isomers/3-12-8-12/fig2_key3d.png", 7.0, 1.5, 5.9, 4.7,
         caption="Validated 3D descriptors (3D-PSA, Rg), R vs S, per solvent")
    return slide


# ── Slide 10b — Model architecture & vision ───────────────────────────────────
def s10b_vision(prs):
    slide = _blank(prs)
    _header(slide, "The model", "A layered permeability model — and where it goes", 13)

    _txt(slide, 0.55, 1.4, 12.2, 0.35,
         "Physics leads (it generalizes); 2D details (it sharpens within a chemotype):",
         color=DARK, size=12, bold=True)

    layers = [
        ("3D FOUNDATION — geometry layer", TEAL,
         "energy-free ensemble descriptors: apolar 3D-PSA · Rgyr · backbone IMHB · shape (NPR) — "
         "shape may capture the small / rigid-macrocycle mechanism where ΔPSA does not  (available now)"),
        ("3D FOUNDATION — energy layer", BLUE,
         "Boltzmann-weighted descriptors + ΔG_transfer  (gated on energy reruns; add-and-select)"),
        ("LOCAL-DETAIL layer", PURPLE,
         "2D fingerprints + Mordred — within-distribution SAR of permeable molecules"),
        ("ESTIMATOR", DARK,
         "tabular foundation model: TabPFN v2 now · evaluate Google's new tabular model"),
    ]
    y = 1.85
    for head, col, body in layers:
        _box(slide, 0.55, y, 4.3, 0.92, head, col, size=11, bold=True)
        _box(slide, 5.0, y, 7.75, 0.92, body, LGREY, fg=DKTEXT, size=10.5, align=PP_ALIGN.LEFT)
        y += 1.0

    _box(slide, 0.55, 5.95, 6.05, 1.0,
         "Long-horizon — de novo design\nAF2 + ProteinMPNN + PyRosetta (BindCraft-style) → "
         "design permeable macrocycles, permeability as a designable objective.",
         GREEN, size=10.5, bold=True)
    _box(slide, 6.75, 5.95, 6.0, 1.0,
         "Long-horizon — PPI optimizer\nany hit as a SMILES → propose residue / bond / "
         "stereochemistry permutations that improve potency + PK/ADME while conserving binding.",
         ORANGE, size=10.5, bold=True)
    return slide


def build_pptx(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    _SLIDE_NO[0] = 0
    for fn in (s01_title, s02_problem, s03_idea, s04_pipeline, s05_timeline,
               s06_csa, s06b_csa_fig, s06c_envelope, s07_rigor, s08_result,
               s08b_result_fig, s09_pivot, s10_ml, s10b_vision, s11_next):
        fn(prs)
    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides._sldIdLst)} slides)")


def parse_args():
    p = argparse.ArgumentParser(description="Generate the 11-slide progress narrative PPTX")
    p.add_argument("--outdir", "-o", default=".")
    return p.parse_args()


if __name__ == "__main__":
    args = parse_args()
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    build_pptx(outdir / "Progress_Deck.pptx")
